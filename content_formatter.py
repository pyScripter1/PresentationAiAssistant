import logging
import json
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE

logger = logging.getLogger(__name__)


class ContentFormatter:
    def __init__(self, design_manager):
        self.design_manager = design_manager

    def format_slide_content(self, slide, slide_info, theme_name):
        """Форматирует контент слайда в зависимости от типа"""
        content_type = slide_info.get('content_type', 'bullet_points')
        content = slide_info.get('content', [])

        if content_type == 'paragraph':
            return self._format_paragraph(slide, content, theme_name)
        elif content_type == 'bullet_points':
            return self._format_bullet_points(slide, content, theme_name)
        elif content_type == 'numbered_list':
            return self._format_numbered_list(slide, content, theme_name)
        elif content_type == 'two_columns':
            return self._format_two_columns(slide, content, theme_name)
        elif content_type == 'title_only':
            return self._format_title_only(slide, theme_name)
        else:
            return self._format_bullet_points(slide, content, theme_name)

    def _format_paragraph(self, slide, content, theme_name):
        """Форматирует как сплошной текст"""
        if not content:
            return

        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        # Объединяем весь контент в один параграф
        full_text = " ".join(content)
        p = text_frame.paragraphs[0]
        p.text = full_text
        p.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY

        # Применяем настройки шрифта
        font_settings = self.design_manager.get_font_settings(theme_name, 'body')
        self._apply_font_settings(p, font_settings)

    def _format_bullet_points(self, slide, content, theme_name):
        """Форматирует как маркированный список"""
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        for i, point in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = f"• {point}"
            p.level = 0

            # Применяем настройки шрифта
            font_settings = self.design_manager.get_font_settings(theme_name, 'bullet_points')
            self._apply_font_settings(p, font_settings)

    def _format_numbered_list(self, slide, content, theme_name):
        """Форматирует как нумерованный список"""
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        for i, point in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = f"{i + 1}. {point}"
            p.level = 0

            font_settings = self.design_manager.get_font_settings(theme_name, 'bullet_points')
            self._apply_font_settings(p, font_settings)

    def _format_two_columns(self, slide, content, theme_name):
        """Форматирует в две колонки - исправленная версия"""
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        # Обрабатываем контент для двух колонок
        left_column_content = []
        right_column_content = []

        # Анализируем структуру контента
        for item in content:
            if isinstance(item, str):
                # Пытаемся распарсить JSON строку
                try:
                    parsed_item = json.loads(item)
                    if isinstance(parsed_item, dict):
                        self._process_column_data(parsed_item, left_column_content, right_column_content)
                    else:
                        # Если это простая строка, добавляем в левую колонку
                        left_column_content.append(item)
                except (json.JSONDecodeError, TypeError):
                    # Если не JSON, добавляем как обычный текст
                    left_column_content.append(item)
            elif isinstance(item, dict):
                self._process_column_data(item, left_column_content, right_column_content)
            else:
                left_column_content.append(str(item))

        # Если не удалось разделить на колонки, делим контент пополам
        if not left_column_content and not right_column_content:
            mid_point = len(content) // 2
            left_column_content = content[:mid_point]
            right_column_content = content[mid_point:]

        # Форматируем левую колонку
        if left_column_content:
            p = text_frame.paragraphs[0]
            p.text = "📘 Основные положения:"
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            font_settings = self.design_manager.get_font_settings(theme_name, 'heading')
            self._apply_font_settings(p, font_settings)

            for point in left_column_content:
                p = text_frame.add_paragraph()
                clean_point = self._clean_content_point(point)
                p.text = f"• {clean_point}"
                p.level = 0
                font_settings = self.design_manager.get_font_settings(theme_name, 'bullet_points')
                self._apply_font_settings(p, font_settings)

        # Добавляем разделитель между колонками
        if left_column_content and right_column_content:
            p = text_frame.add_paragraph()
            p.text = ""  # Пустая строка для разделения

        # Форматируем правую колонку
        if right_column_content:
            p = text_frame.add_paragraph()
            p.text = "📗 Примеры и применение:"
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            font_settings = self.design_manager.get_font_settings(theme_name, 'heading')
            self._apply_font_settings(p, font_settings)

            for point in right_column_content:
                p = text_frame.add_paragraph()
                clean_point = self._clean_content_point(point)
                p.text = f"• {clean_point}"
                p.level = 0
                font_settings = self.design_manager.get_font_settings(theme_name, 'bullet_points')
                self._apply_font_settings(p, font_settings)

    def _process_column_data(self, data, left_column, right_column):
        """Обрабатывает структурированные данные для колонок"""
        if isinstance(data, dict):
            # Обрабатываем структуру с column_title и column_content
            column_title = data.get('column_title', '')
            column_content = data.get('column_content', [])

            if column_title.lower() in ['теория', 'theory', 'основы', 'определения']:
                left_column.append(f"{column_title}:")
                left_column.extend(column_content)
            elif column_title.lower() in ['примеры', 'examples', 'практика', 'применение']:
                right_column.append(f"{column_title}:")
                right_column.extend(column_content)
            else:
                # Если заголовок не распознан, добавляем в левую колонку
                left_column.append(f"{column_title}:")
                left_column.extend(column_content)

    def _clean_content_point(self, point):
        """Очищает пункт контента от лишних символов"""
        if isinstance(point, str):
            # Убираем фигурные скобки и кавычки
            point = point.replace('{', '').replace('}', '').replace('"', '').replace("'", "")
            # Убираем лишние пробелы
            point = ' '.join(point.split())
        return str(point)

    def _format_title_only(self, slide, theme_name):
        """Форматирует слайд только с заголовком"""
        # Очищаем контент (оставляем только заголовок)
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            content_shape.text_frame.clear()

    def _apply_font_settings(self, paragraph, font_settings):
        """Применяет настройки шрифта к параграфу"""
        try:
            for run in paragraph.runs:
                if 'name' in font_settings:
                    run.font.name = font_settings['name']
                if 'size' in font_settings:
                    run.font.size = Pt(font_settings['size'])
                if 'bold' in font_settings:
                    run.font.bold = font_settings['bold']
                if 'color' in font_settings:
                    run.font.color.rgb = self._hex_to_rgb(font_settings['color'])
        except Exception as e:
            logger.error(f"Ошибка применения настроек шрифта: {e}")

    def _hex_to_rgb(self, hex_color):
        """Конвертирует HEX в RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )