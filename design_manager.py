import json
import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.xmlchemy import OxmlElement

logger = logging.getLogger(__name__)


class DesignManager:
    def __init__(self, themes_dir="design_themes"):
        self.themes_dir = themes_dir
        self.available_themes = self._load_available_themes()

    def _load_available_themes(self):
        """Загружает доступные темы оформления"""
        themes = {}
        if os.path.exists(self.themes_dir):
            for file in os.listdir(self.themes_dir):
                if file.endswith('.json'):
                    theme_name = file.replace('.json', '')
                    try:
                        with open(os.path.join(self.themes_dir, file), 'r', encoding='utf-8') as f:
                            themes[theme_name] = json.load(f)
                    except Exception as e:
                        logger.error(f"Ошибка загрузки темы {file}: {e}")
        return themes

    def get_available_themes(self):
        """Возвращает список доступных тем"""
        return list(self.available_themes.keys())

    def apply_theme_to_presentation(self, prs, theme_name):
        """Применяет тему ко всей презентации"""
        theme = self.available_themes.get(theme_name)
        if not theme:
            logger.warning(f"Тема {theme_name} не найдена, используется стандартная")
            return prs

        # Применяем настройки темы
        self._apply_background(prs, theme)
        return prs

    def _apply_background(self, prs, theme):
        """Применяет фон к слайдам"""
        background_config = theme.get('background', {})
        bg_type = background_config.get('type', 'solid')

        for slide in prs.slides:
            if bg_type == 'gradient':
                self._set_gradient_background(slide, background_config)
            else:
                self._set_solid_background(slide, background_config)

    def _set_solid_background(self, slide, background_config):
        """Устанавливает сплошной фон"""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self._hex_to_rgb(background_config.get('color1', '#FFFFFF'))
        except Exception as e:
            logger.error(f"Ошибка установки фона: {e}")

    def _set_gradient_background(self, slide, background_config):
        """Устанавливает градиентный фон"""
        try:
            background = slide.background
            fill = background.fill
            fill.gradient()
            # Настройки градиента можно добавить здесь
        except Exception as e:
            logger.error(f"Ошибка установки градиентного фона: {e}")

    def apply_slide_design(self, slide, slide_type, theme_name):
        """Применяет дизайн к конкретному слайду"""
        theme = self.available_themes.get(theme_name)
        if not theme:
            return

        layouts = theme.get('layouts', {})
        slide_layout = layouts.get(slide_type, layouts.get('content_slide', {}))

        # Применяем цвета к заголовку
        if slide.shapes.title:
            title_color = slide_layout.get('title_color')
            if title_color:
                self._set_text_color(slide.shapes.title, title_color)

    def _set_text_color(self, shape, hex_color):
        """Устанавливает цвет текста"""
        try:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self._hex_to_rgb(hex_color)
        except Exception as e:
            logger.error(f"Ошибка установки цвета текста: {e}")

    def _hex_to_rgb(self, hex_color):
        """Конвертирует HEX в RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def get_font_settings(self, theme_name, font_type):
        """Возвращает настройки шрифта для темы"""
        theme = self.available_themes.get(theme_name)
        if not theme:
            return {}
        return theme.get('fonts', {}).get(font_type, {})