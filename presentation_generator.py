import requests
import json
import os
import re
import logging
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from config import Config
from design_manager import DesignManager
from content_formatter import ContentFormatter

# Настройка логирования
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PresentationGenerator:
    def __init__(self):
        self.api_key = Config.MISTRAL_API_KEY
        self.presentations_dir = Config.PRESENTATIONS_DIR
        self.templates_dir = "templates"

        # Инициализация менеджеров
        self.design_manager = DesignManager()
        self.content_formatter = ContentFormatter(self.design_manager)

        # Создаем папки, если они не существуют
        os.makedirs(self.presentations_dir, exist_ok=True)
        os.makedirs(self.templates_dir, exist_ok=True)
        os.makedirs("design_themes", exist_ok=True)

        # Загружаем доступные шаблоны
        self.available_templates = self._load_available_templates()

        # Создаем дефолтный шаблон если нет шаблонов
        if not self.available_templates:
            self._create_default_template()

        # Создаем дефолтные темы если нет тем
        if not self.design_manager.get_available_themes():
            self._create_default_themes()

    def _load_available_templates(self):
        """Загружает список доступных шаблонов"""
        templates = {}
        if os.path.exists(self.templates_dir):
            for file in os.listdir(self.templates_dir):
                if file.endswith('.json'):
                    template_name = file.replace('.json', '')
                    try:
                        with open(os.path.join(self.templates_dir, file), 'r', encoding='utf-8') as f:
                            template_data = json.load(f)
                            templates[template_name] = template_data
                            logger.info(f"Загружен шаблон: {template_name}")
                    except Exception as e:
                        logger.error(f"Ошибка загрузки шаблона {file}: {e}")
        return templates

    def _create_default_templates(self):
        """Создает все дефолтные шаблоны"""
        templates = {
            "educational": {
                "name": "educational",
                "description": "Образовательный шаблон с разными типами контента",
                "default_design_theme": "modern_blue",
                "slides_structure": [
                    {
                        "slide_number": 1,
                        "slide_type": "title",
                        "content_type": "title_only",
                        "slide_title": "Название темы",
                        "content": [],
                        "instructions": "Создай заголовок презентации на указанную тему"
                    },
                    {
                        "slide_number": 2,
                        "slide_type": "objectives",
                        "content_type": "bullet_points",
                        "slide_title": "Цели презентации",
                        "content": [],
                        "instructions": "Опиши 3-4 основные цели этой презентации в виде пунктов"
                    },
                    {
                        "slide_number": 3,
                        "slide_type": "introduction",
                        "content_type": "paragraph",
                        "slide_title": "Введение и теоретические определения",
                        "content": [],
                        "instructions": "Напиши связный текст на 150-200 слов с введением в тему"
                    },
                    {
                        "slide_number": 4,
                        "slide_type": "main_content",
                        "content_type": "two_columns",
                        "slide_title": "Основная часть",
                        "content": [],
                        "instructions": "Представь информацию в двух колонках: теория слева, примеры справа"
                    },
                    {
                        "slide_number": 5,
                        "slide_type": "examples",
                        "content_type": "bullet_points",
                        "slide_title": "Примеры и практическое применение",
                        "content": [],
                        "instructions": "Приведи 3-4 конкретных примера в виде пунктов"
                    },
                    {
                        "slide_number": 6,
                        "slide_type": "exercises",
                        "content_type": "numbered_list",
                        "slide_title": "Задания для самостоятельного решения",
                        "content": [],
                        "instructions": "Создай 3-4 задания с нумерацией для самостоятельной работы"
                    },
                    {
                        "slide_number": 7,
                        "slide_type": "conclusion",
                        "content_type": "paragraph",
                        "slide_title": "Заключение и выводы",
                        "content": [],
                        "instructions": "Напиши заключительный текст на 120-180 слов с выводами"
                    },
                    {
                        "slide_number": 8,
                        "slide_type": "homework",
                        "content_type": "bullet_points",
                        "slide_title": "Домашнее задание",
                        "content": [],
                        "instructions": "Предложи 3-4 пункта домашнего задания"
                    }
                ]
            },
            "corporate": {
                "name": "corporate",
                "description": "Корпоративный шаблон для бизнес-презентаций",
                "default_design_theme": "modern_blue",
                "slides_structure": [
                    {
                        "slide_number": 1,
                        "slide_type": "title",
                        "content_type": "title_only",
                        "slide_title": "Название проекта/презентации",
                        "content": [],
                        "instructions": "Создай профессиональный заголовок для бизнес-презентации"
                    },
                    {
                        "slide_number": 2,
                        "slide_type": "agenda",
                        "content_type": "bullet_points",
                        "slide_title": "Повестка встречи",
                        "content": [],
                        "instructions": "Создай 4-5 пунктов повестки встречи в деловом стиле"
                    },
                    {
                        "slide_number": 3,
                        "slide_type": "executive_summary",
                        "content_type": "paragraph",
                        "slide_title": "Краткий обзор",
                        "content": [],
                        "instructions": "Напиши краткий обзор на 150-200 слов в деловом стиле"
                    },
                    {
                        "slide_number": 4,
                        "slide_type": "problem_statement",
                        "content_type": "bullet_points",
                        "slide_title": "Постановка проблемы",
                        "content": [],
                        "instructions": "Опиши 3-4 ключевые проблемы или вызовы"
                    },
                    {
                        "slide_number": 5,
                        "slide_type": "solution",
                        "content_type": "two_columns",
                        "slide_title": "Предлагаемое решение",
                        "content": [],
                        "instructions": "Опиши решение: левая колонка - подход, правая - преимущества"
                    },
                    {
                        "slide_number": 6,
                        "slide_type": "metrics",
                        "content_type": "bullet_points",
                        "slide_title": "Ключевые метрики",
                        "content": [],
                        "instructions": "Представь 4-5 измеримых показателей успеха"
                    },
                    {
                        "slide_number": 7,
                        "slide_type": "timeline",
                        "content_type": "numbered_list",
                        "slide_title": "План реализации",
                        "content": [],
                        "instructions": "Создай нумерованный план из 4-5 этапов с сроками"
                    },
                    {
                        "slide_number": 8,
                        "slide_type": "team",
                        "content_type": "bullet_points",
                        "slide_title": "Команда проекта",
                        "content": [],
                        "instructions": "Опиши ключевых участников команды и их роли"
                    },
                    {
                        "slide_number": 9,
                        "slide_type": "budget",
                        "content_type": "bullet_points",
                        "slide_title": "Бюджет и ресурсы",
                        "content": [],
                        "instructions": "Опиши основные статьи бюджета и необходимые ресурсы"
                    },
                    {
                        "slide_number": 10,
                        "slide_type": "next_steps",
                        "content_type": "numbered_list",
                        "slide_title": "Следующие шаги",
                        "content": [],
                        "instructions": "Определи 3-4 конкретных следующих шага с ответственными"
                    }
                ]
            },
            "creative": {
                "name": "creative",
                "description": "Креативный шаблон для презентаций и выступлений",
                "default_design_theme": "elegant_green",
                "slides_structure": [
                    {
                        "slide_number": 1,
                        "slide_type": "title",
                        "content_type": "title_only",
                        "slide_title": "Захватывающий заголовок",
                        "content": [],
                        "instructions": "Создай креативный и запоминающийся заголовок"
                    },
                    {
                        "slide_number": 2,
                        "slide_type": "hook",
                        "content_type": "paragraph",
                        "slide_title": "История, которая зацепит",
                        "content": [],
                        "instructions": "Напиши короткую историю или пример, который привлечет внимание аудитории"
                    },
                    {
                        "slide_number": 3,
                        "slide_type": "big_idea",
                        "content_type": "paragraph",
                        "slide_title": "Основная идея",
                        "content": [],
                        "instructions": "Сформулируй основную идею презентации вдохновляющим языком"
                    },
                    {
                        "slide_number": 4,
                        "slide_type": "why_matters",
                        "content_type": "bullet_points",
                        "slide_title": "Почему это важно?",
                        "content": [],
                        "instructions": "Объясни 3-4 причины, почему эта тема важна для аудитории"
                    },
                    {
                        "slide_number": 5,
                        "slide_type": "surprising_facts",
                        "content_type": "bullet_points",
                        "slide_title": "Неожиданные факты",
                        "content": [],
                        "instructions": "Представь 3-4 удивительных или неожиданных факта по теме"
                    },
                    {
                        "slide_number": 6,
                        "slide_type": "visual_story",
                        "content_type": "paragraph",
                        "slide_title": "Визуальная история",
                        "content": [],
                        "instructions": "Опиши визуальный пример или кейс, который иллюстрирует тему"
                    },
                    {
                        "slide_number": 7,
                        "slide_type": "how_it_works",
                        "content_type": "two_columns",
                        "slide_title": "Как это работает?",
                        "content": [],
                        "instructions": "Объясни механизм работы: левая колонка - принципы, правая - примеры"
                    },
                    {
                        "slide_number": 8,
                        "slide_type": "success_stories",
                        "content_type": "bullet_points",
                        "slide_title": "Истории успеха",
                        "content": [],
                        "instructions": "Расскажи 3-4 вдохновляющие истории успеха или кейсы"
                    },
                    {
                        "slide_number": 9,
                        "slide_type": "call_to_action",
                        "content_type": "numbered_list",
                        "slide_title": "Призыв к действию",
                        "content": [],
                        "instructions": "Создай 3-4 конкретных шага, которые аудитория может сделать прямо сейчас"
                    },
                    {
                        "slide_number": 10,
                        "slide_type": "takeaways",
                        "content_type": "bullet_points",
                        "slide_title": "Ключевые выводы",
                        "content": [],
                        "instructions": "Сформулируй 3-4 главных вывода из презентации"
                    }
                ]
            }
        }

        for template_name, template_data in templates.items():
            try:
                template_path = os.path.join(self.templates_dir, f"{template_name}.json")
                with open(template_path, 'w', encoding='utf-8') as f:
                    json.dump(template_data, f, ensure_ascii=False, indent=2)
                self.available_templates[template_name] = template_data
                logger.info(f"Создан шаблон: {template_name}")
            except Exception as e:
                logger.error(f"Ошибка создания шаблона {template_name}: {e}")

    def _create_default_themes(self):
        """Создает дефолтные темы оформления"""
        themes = {
            "modern_blue": {
                "name": "modern_blue",
                "description": "Современная синяя тема для образовательных презентаций",
                "colors": {
                    "primary": "#2C5AA0",
                    "secondary": "#4A90E2",
                    "accent": "#FF6B6B",
                    "background": "#FFFFFF",
                    "text_primary": "#2C3E50",
                    "text_secondary": "#7F8C8D"
                },
                "fonts": {
                    "title": {"name": "Calibri", "size": 44, "bold": True, "color": "#2C5AA0"},
                    "heading": {"name": "Calibri", "size": 32, "bold": True, "color": "#2C3E50"},
                    "body": {"name": "Calibri", "size": 18, "bold": False, "color": "#2C3E50"},
                    "bullet_points": {"name": "Calibri", "size": 16, "bold": False, "color": "#2C3E50"}
                },
                "background": {
                    "type": "solid",
                    "color1": "#FFFFFF"
                },
                "layouts": {
                    "title_slide": {
                        "background_color": "#2C5AA0",
                        "title_color": "#FFFFFF"
                    },
                    "content_slide": {
                        "background_color": "#FFFFFF",
                        "title_color": "#2C5AA0",
                        "content_color": "#2C3E50"
                    }
                }
            },
            "elegant_green": {
                "name": "elegant_green",
                "description": "Элегантная зеленая тема",
                "colors": {
                    "primary": "#27AE60",
                    "secondary": "#2ECC71",
                    "accent": "#E74C3C",
                    "background": "#FDFEFE",
                    "text_primary": "#2C3E50",
                    "text_secondary": "#566573"
                },
                "fonts": {
                    "title": {"name": "Calibri", "size": 44, "bold": True, "color": "#27AE60"},
                    "heading": {"name": "Calibri", "size": 32, "bold": True, "color": "#2C3E50"},
                    "body": {"name": "Calibri", "size": 18, "bold": False, "color": "#2C3E50"},
                    "bullet_points": {"name": "Calibri", "size": 16, "bold": False, "color": "#2C3E50"}
                },
                "background": {
                    "type": "solid",
                    "color1": "#FDFEFE"
                },
                "layouts": {
                    "title_slide": {
                        "background_color": "#27AE60",
                        "title_color": "#FFFFFF"
                    },
                    "content_slide": {
                        "background_color": "#FDFEFE",
                        "title_color": "#27AE60",
                        "content_color": "#2C3E50"
                    }
                }
            }
        }

        for theme_name, theme_data in themes.items():
            try:
                theme_path = os.path.join("design_themes", f"{theme_name}.json")
                with open(theme_path, 'w', encoding='utf-8') as f:
                    json.dump(theme_data, f, ensure_ascii=False, indent=2)
                logger.info(f"Создана тема: {theme_name}")
            except Exception as e:
                logger.error(f"Ошибка создания темы {theme_name}: {e}")

    def get_available_templates(self):
        """Возвращает список доступных шаблонов"""
        return list(self.available_templates.keys())

    def get_available_themes(self):
        """Возвращает список доступных тем оформления"""
        return self.design_manager.get_available_themes()

    def load_template(self, template_name):
        """Загружает конкретный шаблон"""
        template = self.available_templates.get(template_name)
        if not template:
            raise ValueError(
                f"Шаблон '{template_name}' не найден. Доступные шаблоны: {list(self.available_templates.keys())}")
        return template

    def query_mistral_api(self, prompt, model="mistral-small-latest"):
        """Отправляет запрос к Mistral AI API и возвращает ответ."""
        url = "https://api.mistral.ai/v1/chat/completions"

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }

        data = {
            "model": model,
            "messages": [
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7
        }

        try:
            logger.info("Отправка запроса к Mistral AI...")
            response = requests.post(url, headers=headers, data=json.dumps(data), timeout=60)
            response.raise_for_status()
            response_data = response.json()
            logger.info("Успешный ответ от Mistral AI")
            return response_data["choices"][0]["message"]["content"]

        except requests.exceptions.Timeout:
            logger.error("Таймаут при запросе к Mistral AI")
            return None
        except requests.exceptions.RequestException as e:
            logger.error(f"Ошибка при запросе к API: {e}")
            return None
        except KeyError as e:
            logger.error(f"Неожиданный формат ответа от API: {e}")
            return None
        except Exception as e:
            logger.error(f"Неожиданная ошибка: {e}")
            return None

    def create_presentation_prompt(self, topic, template_name="educational", additional_prompt=""):
        """Создает промпт для генерации презентации по шаблону."""

        template = self.load_template(template_name)
        if not template:
            raise ValueError(f"Шаблон '{template_name}' не найден")

        prompt = f"""
    Ты — опытный преподаватель и создатель образовательных материалов. 
    Твоя задача — создать подробную структуру и контент для учебной презентации на тему "{topic}".

    СТРУКТУРА ПРЕЗЕНТАЦИИ:
    """

        # Добавляем инструкции для каждого слайда из шаблона
        for slide_template in template["slides_structure"]:
            content_type_instruction = ""
            if slide_template["content_type"] == "paragraph":
                content_type_instruction = "Напиши связный текст без пунктов. Объедини все в один абзац."
            elif slide_template["content_type"] == "bullet_points":
                content_type_instruction = "Используй маркированный список с 3-4 четкими пунктами."
            elif slide_template["content_type"] == "numbered_list":
                content_type_instruction = "Используй нумерованный список с 3-4 пунктами для заданий."
            elif slide_template["content_type"] == "two_columns":
                content_type_instruction = "Для двух колонок создай 4-6 общих пунктов. Я сам разделю их на колонки. НЕ используй JSON структуры!"
            elif slide_template["content_type"] == "title_only":
                content_type_instruction = "Только заголовок, без дополнительного контента."

            prompt += f"""
    Слайд {slide_template['slide_number']}: {slide_template['slide_title']}
    Тип контента: {content_type_instruction}
    Задача: {slide_template['instructions']}
    """

        # Добавляем дополнительные пожелания
        if additional_prompt:
            prompt += f"\nДОПОЛНИТЕЛЬНЫЕ ТРЕБОВАНИЯ:\n{additional_prompt}\n"

        prompt += f"""
    ВЕРНИ ОТВЕТ В ФОРМАТЕ JSON:

    {{
      "presentation_title": "Название презентации",
      "template_used": "{template_name}",
      "slides": [
        {{
          "slide_number": 1,
          "slide_type": "title",
          "content_type": "title_only",
          "slide_title": "Название темы",
          "content": []
        }},
        {{
          "slide_number": 2,
          "slide_type": "objectives", 
          "content_type": "bullet_points",
          "slide_title": "Цели презентации",
          "content": [
            "Понять основные концепции темы",
            "Научиться применять знания на практике",
            "Развить критическое мышление"
          ]
        }},
        {{
          "slide_number": 3,
          "slide_type": "introduction",
          "content_type": "paragraph", 
          "slide_title": "Введение",
          "content": [
            "Здесь будет связный текст введения без пунктов..."
          ]
        }},
        {{
          "slide_number": 4,
          "slide_type": "main_content",
          "content_type": "two_columns",
          "slide_title": "Основная часть",
          "content": [
            "Основное понятие 1 с объяснением",
            "Основное понятие 2 с характеристиками", 
            "Практический пример 1 для иллюстрации",
            "Практический пример 2 из реальной жизни",
            "Ключевой вывод по теме",
            "Дополнительная информация для понимания"
          ]
        }}
        // ... добавь остальные слайды по структуре шаблона
      ]
    }}

    ВАЖНЫЕ ПРАВИЛА:
    1. Для two_columns НЕ используй JSON структуры - только простой текст в массиве content
    2. Для paragraph объединяй весь текст в один элемент массива
    3. Для bullet_points и numbered_list используй массив с отдельными пунктами
    4. Всегда экранируй кавычки в тексте

    Тема: {topic}
    """
        return prompt

    def create_pptx_from_data(self, presentation_data, output_filename, design_theme=None):
        """Создает PowerPoint файл с применением дизайна"""
        try:
            logger.info(f"Создание PPTX файла с дизайном: {output_filename}")
            prs = Presentation()

            # Настройки презентации
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Применяем тему дизайна
            theme_name = design_theme or presentation_data.get('design_theme', 'modern_blue')
            self.design_manager.apply_theme_to_presentation(prs, theme_name)

            # Создаем слайды
            for slide_info in presentation_data["slides"]:
                self._create_slide(prs, slide_info, theme_name)

            prs.save(output_filename)
            logger.info(f"PPTX файл успешно создан: {output_filename}")
            return output_filename

        except Exception as e:
            logger.error(f"Ошибка при создании PPTX файла: {e}")
            return None

    def _create_slide(self, prs, slide_info, theme_name):
        """Создает один слайд с применением дизайна и форматирования"""
        slide_type = slide_info.get("slide_type", "content")
        content_type = slide_info.get("content_type", "bullet_points")

        # Выбираем макет в зависимости от типа слайда
        if slide_type == "title" or content_type == "title_only":
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Титульный
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Контентный

        # Устанавливаем заголовок
        title_shape = slide.shapes.title
        if slide_info["slide_title"]:
            title_shape.text = slide_info["slide_title"]

        # Применяем дизайн к слайду
        self.design_manager.apply_slide_design(slide, slide_type, theme_name)

        # Форматируем контент (если это не титульный слайд)
        if content_type != "title_only" and len(slide.placeholders) > 1:
            self.content_formatter.format_slide_content(slide, slide_info, theme_name)

    def _fix_json_errors(self, json_text):
        """Исправляет распространенные ошибки в JSON от ИИ"""
        try:
            # Убираем Markdown разметку если есть
            cleaned = json_text.strip()
            if cleaned.startswith("```json"):
                cleaned = cleaned[7:]
            elif cleaned.startswith("```"):
                cleaned = cleaned[3:]
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]

            # Исправляем неэкранированные кавычки внутри строк
            lines = cleaned.split('\n')
            fixed_lines = []
            in_string = False
            escape_next = False

            for line in lines:
                new_line = ""
                for char in line:
                    if char == '"' and not escape_next:
                        in_string = not in_string
                        new_line += char
                    elif char == '\\' and in_string:
                        escape_next = True
                        new_line += char
                    else:
                        if escape_next:
                            escape_next = False
                        new_line += char
                fixed_lines.append(new_line)

            cleaned = '\n'.join(fixed_lines)

            # Убираем лишние запятые
            cleaned = re.sub(r',\s*}', '}', cleaned)
            cleaned = re.sub(r',\s*]', ']', cleaned)

            return cleaned

        except Exception as e:
            logger.error(f"Ошибка при исправлении JSON: {e}")
            return json_text

    def clean_json_response(self, response_text):
        """Очищает ответ от модели от лишних символов и извлекает JSON."""
        try:
            if not response_text:
                raise ValueError("Пустой ответ от API")

            # Сначала пытаемся распарсить как есть
            try:
                presentation_data = json.loads(response_text)
                logger.info("JSON успешно обработан без исправлений")
                return presentation_data
            except json.JSONDecodeError:
                pass

            # Если не получилось, очищаем и пытаемся снова
            cleaned_response = self._fix_json_errors(response_text)

            try:
                presentation_data = json.loads(cleaned_response)
                logger.info("JSON успешно обработан после исправлений")
                return presentation_data
            except json.JSONDecodeError as e:
                # Последняя попытка - найти JSON в тексте
                logger.info("Поиск JSON pattern в тексте...")
                json_match = re.search(r'\{.*\}', cleaned_response, re.DOTALL)
                if json_match:
                    json_str = json_match.group(0)
                    try:
                        presentation_data = json.loads(json_str)
                        logger.info("JSON найден и обработан через regex")
                        return presentation_data
                    except json.JSONDecodeError:
                        pass

                logger.error(f"Не удалось распарсить JSON после всех попыток: {e}")
                logger.error(f"Очищенный текст: {cleaned_response[:500]}...")
                return None

        except Exception as e:
            logger.error(f"Ошибка при обработке ответа: {e}")
            return None

    def generate_presentation(self, topic, template_name="educational", additional_prompt="", design_theme=None):
        """Основной метод для генерации презентации по шаблону."""

        try:
            # Проверяем входные параметры
            if not topic or not topic.strip():
                raise ValueError("Тема презентации не может быть пустой")

            topic = topic.strip()

            # Генерируем имя файла
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_topic = "".join(c if c.isalnum() else "_" for c in topic)
            output_filename = os.path.join(
                self.presentations_dir,
                f"presentation_{safe_topic}_{timestamp}.pptx"
            )

            logger.info(f"Начало генерации презентации: '{topic}', шаблон: '{template_name}', дизайн: '{design_theme}'")

            # Создаем промпт и получаем ответ от Mistral AI
            prompt = self.create_presentation_prompt(topic, template_name, additional_prompt)
            response_text = self.query_mistral_api(prompt)

            if not response_text:
                raise Exception("Не удалось получить ответ от Mistral AI")

            logger.info("Ответ от ИИ получен, обрабатываю...")

            presentation_data = self.clean_json_response(response_text)

            if not presentation_data:
                raise Exception("Не удалось обработать ответ от ИИ")

            # Если тема не указана, используем случайную из доступных
            if not design_theme:
                available_themes = self.get_available_themes()
                design_theme = random.choice(available_themes) if available_themes else 'modern_blue'

            result_file = self.create_pptx_from_data(presentation_data, output_filename, design_theme)

            if not result_file:
                raise Exception("Ошибка при создании PPTX файла")

            logger.info(f"Презентация успешно создана: {result_file}")
            return result_file

        except Exception as e:
            logger.error(f"Ошибка при генерации презентации: {e}")
            return None

    def get_presentation_info(self, presentation_data):
        """Возвращает информацию о презентации для отправки пользователю."""
        if not presentation_data:
            return None

        info = {
            'title': presentation_data.get('presentation_title', 'Без названия'),
            'template': presentation_data.get('template_used', 'unknown'),
            'slides_count': len(presentation_data.get('slides', [])),
            'structure': []
        }

        for slide in presentation_data.get('slides', []):
            info['structure'].append({
                'number': slide.get('slide_number'),
                'title': slide.get('slide_title'),
                'type': slide.get('slide_type', 'content'),
                'content_type': slide.get('content_type', 'bullet_points')
            })

        return info


# Функция для обратной совместимости
def generate_presentation(topic, template_name="educational", additional_prompt="", design_theme=None):
    generator = PresentationGenerator()
    return generator.generate_presentation(topic, template_name, additional_prompt, design_theme)


if __name__ == "__main__":
    # Тестовый запуск с подробным логированием
    print("=== Тест генератора презентаций ===")
    generator = PresentationGenerator()

    print(f"Доступные шаблоны: {generator.get_available_templates()}")
    print(f"Доступные темы: {generator.get_available_themes()}")
    print(f"API ключ: {'установлен' if Config.MISTRAL_API_KEY else 'отсутствует'}")

    topic = input("Введите тему для теста: ").strip()

    result = generator.generate_presentation(
        topic=topic,
        template_name="educational",
        additional_prompt="Количество слайдов: 8",
        design_theme="modern_blue"
    )

    if result:
        print(f"✅ Тест успешен! Файл создан: {result}")
    else:
        print("❌ Тест не удался. Проверьте логи для деталей.")